from __future__ import annotations

import sys
from pathlib import Path

REPO_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(REPO_ROOT / "tmp" / "presentation_py"))

from PIL import Image  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE  # noqa: E402
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

from build_presentation import ASSET_DIR, PROJECT_UNDERSTANDING, SCENARIOS  # noqa: E402


OUT_DIR = REPO_ROOT / "presentation"
OUT_FILE = OUT_DIR / "indeng243_diet_planner_presentation_editable.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG = "F6F1E8"
GREEN = "204736"
GREEN_2 = "315B46"
TEXT = "163126"
MUTED = "5D6E61"
CARD = "FFFDF8"
LINE = "D8D2C6"
SOFT_GREEN = "EAF2E7"
SOFT_BLUE = "E8F0F7"
SOFT_TERRA = "F8E8E1"
SOFT_GOLD = "F5E5BE"
ACCENT = "D9A441"

FONT_TITLE = "Avenir Next"
FONT_BODY = "Avenir"
FONT_BOLD = "Avenir Next"


def rgb(hex_value: str) -> RGBColor:
    return RGBColor.from_string(hex_value)


def add_full_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(BG)


def add_top_bar(slide, kicker: str, slide_number: int) -> None:
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(GREEN)
    bar.line.fill.background()

    kicker_box = slide.shapes.add_textbox(Inches(0.58), Inches(0.12), Inches(2.6), Inches(0.24))
    tf = kicker_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = kicker.upper()
    p.font.name = FONT_BOLD
    p.font.size = Pt(15)
    p.font.color.rgb = rgb(ACCENT)

    num_box = slide.shapes.add_textbox(Inches(12.6), Inches(7.05), Inches(0.35), Inches(0.2))
    tf = num_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    p.text = f"{slide_number:02d}"
    p.font.name = FONT_BOLD
    p.font.size = Pt(15)
    p.font.color.rgb = rgb(MUTED)


def style_text_frame(shape, *, left=8, right=8, top=6, bottom=6, anchor=MSO_ANCHOR.TOP) -> None:
    tf = shape.text_frame
    tf.margin_left = Pt(left)
    tf.margin_right = Pt(right)
    tf.margin_top = Pt(top)
    tf.margin_bottom = Pt(bottom)
    tf.word_wrap = True
    tf.vertical_anchor = anchor


def add_title(slide, title: str, subtitle: str, *, top=0.86) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.58), Inches(top), Inches(10.9), Inches(0.75))
    style_text_frame(title_box, left=0, right=0, top=0, bottom=0)
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = FONT_TITLE
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = rgb(TEXT)

    sub_box = slide.shapes.add_textbox(Inches(0.58), Inches(top + 0.56), Inches(10.8), Inches(0.55))
    style_text_frame(sub_box, left=0, right=0, top=0, bottom=0)
    p = sub_box.text_frame.paragraphs[0]
    p.text = subtitle
    p.font.name = FONT_BODY
    p.font.size = Pt(16)
    p.font.color.rgb = rgb(MUTED)


def add_card(slide, left, top, width, height, *, fill= CARD, line=LINE, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(1)
    return shape


def add_textbox(slide, left, top, width, height, text: str, *, font_size=18, bold=False, color=TEXT, font_name=FONT_BODY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    style_text_frame(box)
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)
    return box


def add_bullets(slide, left, top, width, height, bullets: list[str], *, font_size=17, color=TEXT, line_space=1.2):
    box = slide.shapes.add_textbox(left, top, width, height)
    style_text_frame(box)
    tf = box.text_frame
    tf.clear()
    first = True
    for bullet in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = f"• {bullet}"
        p.font.name = FONT_BODY
        p.font.size = Pt(font_size)
        p.font.color.rgb = rgb(color)
        p.line_spacing = line_space
        p.space_after = Pt(6)
    return box


def add_labeled_metric(slide, left, top, width, height, label: str, value: str, *, fill=SOFT_GREEN):
    add_card(slide, left, top, width, height, fill=fill, line=fill)
    add_textbox(slide, left + Inches(0.12), top + Inches(0.1), Inches(1.0), Inches(0.24), label, font_size=15, bold=True, color=GREEN, font_name=FONT_BOLD)
    add_textbox(slide, left + Inches(1.25), top + Inches(0.1), width - Inches(1.35), Inches(0.24), value, font_size=15, color=TEXT)


def add_tag(slide, left, top, text: str, *, fill=SOFT_GREEN, color=GREEN):
    tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, Inches(2.1), Inches(0.33))
    tag.fill.solid()
    tag.fill.fore_color.rgb = rgb(fill)
    tag.line.fill.background()
    tf = tag.text_frame
    style_text_frame(tag, left=8, right=8, top=4, bottom=2, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT_BOLD
    p.font.size = Pt(14)
    p.font.color.rgb = rgb(color)
    return tag


def image_size(path: Path) -> tuple[int, int]:
    with Image.open(path) as img:
        return img.width, img.height


def add_picture_contain(slide, path: Path, left, top, width, height):
    img_w, img_h = image_size(path)
    img_ratio = img_w / img_h
    box_ratio = width / height
    if img_ratio > box_ratio:
        pic_w = width
        pic_h = int(width / img_ratio)
        pic_left = left
        pic_top = top + int((height - pic_h) / 2)
    else:
        pic_h = height
        pic_w = int(height * img_ratio)
        pic_top = top
        pic_left = left + int((width - pic_w) / 2)
    slide.shapes.add_picture(str(path), pic_left, pic_top, width=pic_w, height=pic_h)


def add_image_frame(slide, left, top, width, height, path: Path):
    frame = add_card(slide, left, top, width, height, fill=CARD, line=LINE)
    frame.shadow.inherit = False
    add_picture_contain(slide, path, left + Inches(0.04), top + Inches(0.04), width - Inches(0.08), height - Inches(0.08))


def build_slide_1(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide)
    add_top_bar(slide, "INDENG 243 Analytics Communication Final", 1)

    circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(9.9), Inches(-0.1), Inches(3.5), Inches(3.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = rgb(GREEN_2)
    circle.line.fill.background()

    add_textbox(slide, Inches(0.58), Inches(1.15), Inches(4.8), Inches(1.35), "From Nutrition Goals\nto a Grocery Trip", font_size=36, bold=True, font_name=FONT_TITLE)
    add_textbox(
        slide,
        Inches(0.58),
        Inches(2.64),
        Inches(5.2),
        Inches(1.0),
        "This tool turns a health goal into a practical grocery list, a rough cost estimate, and nearby store suggestions.",
        font_size=18,
        color=MUTED,
    )
    add_tag(slide, Inches(0.58), Inches(3.82), "Non-technical audience first", fill=SOFT_GREEN, color=GREEN)
    add_tag(slide, Inches(0.58), Inches(4.22), "Observed on the deployed app", fill=SOFT_BLUE, color="245072")
    add_tag(slide, Inches(0.58), Inches(4.62), "Focus: decisions, tradeoffs, outcomes", fill=SOFT_TERRA, color="C86A45")

    add_card(slide, Inches(0.58), Inches(5.42), Inches(4.45), Inches(1.15), fill=CARD, line=LINE)
    add_textbox(slide, Inches(0.72), Inches(5.58), Inches(4.1), Inches(0.22), "What the audience sees", font_size=16, bold=True, color=GREEN, font_name=FONT_BOLD)
    add_bullets(
        slide,
        Inches(0.72),
        Inches(5.88),
        Inches(4.0),
        Inches(0.88),
        [
            "Preset goals like muscle gain, fat loss, maintenance, and budget-friendly healthy.",
            "A one-click path from location and targets to a shopping list and store suggestions.",
            "Outputs that read like a grocery assistant, not a technical dashboard.",
        ],
        font_size=13,
        color=TEXT,
    )

    add_image_frame(slide, Inches(7.05), Inches(0.95), Inches(5.55), Inches(5.45), ASSET_DIR / "dashboard-home.png")
    add_textbox(slide, Inches(7.05), Inches(6.48), Inches(5.45), Inches(0.3), "Hero view: the deployed dashboard starts with presets, local store lookup, and a guided form.", font_size=13, color=MUTED)


def build_slide_2(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide)
    add_top_bar(slide, "Audience", 2)
    add_title(
        slide,
        "Who This Is For",
        "The clearest non-technical audience is a shopper who wants healthier and easier grocery decisions, not a model developer.",
    )

    add_card(slide, Inches(0.58), Inches(2.0), Inches(4.35), Inches(4.35), fill=CARD, line=LINE)
    add_textbox(slide, Inches(0.78), Inches(2.2), Inches(3.8), Inches(0.26), "Primary user persona", font_size=17, bold=True, color=GREEN, font_name=FONT_BOLD)
    add_textbox(slide, Inches(0.78), Inches(2.6), Inches(3.9), Inches(0.55), "Busy students and young professionals", font_size=24, bold=True, font_name=FONT_TITLE)
    add_textbox(
        slide,
        Inches(0.78),
        Inches(3.3),
        Inches(3.9),
        Inches(1.25),
        "People who want to shop for muscle gain, fat loss, maintenance, or a healthier budget without building a meal plan from scratch.",
        font_size=17,
        color=MUTED,
    )
    add_bullets(
        slide,
        Inches(0.78),
        Inches(4.7),
        Inches(3.95),
        Inches(1.4),
        [
            "They want a recommendation in minutes, not a spreadsheet exercise.",
            "They care about rough basket cost before they get to the store.",
            "They need outputs they can trust without knowing the math.",
        ],
        font_size=14,
    )

    boxes = [
        (Inches(5.25), Inches(2.0), "Fast setup", "Preset buttons and a one-click recommend flow lower setup effort.", SOFT_GREEN),
        (Inches(9.2), Inches(2.0), "Affordable", "The app surfaces a typical basket cost so the recommendation feels practical.", SOFT_BLUE),
        (Inches(5.25), Inches(4.15), "Local", "Nearby store lookup makes the output feel tied to a real trip, not just a nutrition target.", SOFT_TERRA),
        (Inches(9.2), Inches(4.15), "Understandable", "The interface uses grocery language and simple controls instead of technical settings.", SOFT_GOLD),
    ]
    for left, top, title, body, fill in boxes:
        add_card(slide, left, top, Inches(3.45), Inches(1.7), fill=fill, line=fill)
        add_textbox(slide, left + Inches(0.18), top + Inches(0.18), Inches(3.0), Inches(0.22), title, font_size=18, bold=True, color=GREEN, font_name=FONT_BOLD)
        add_textbox(slide, left + Inches(0.18), top + Inches(0.58), Inches(3.05), Inches(0.9), body, font_size=15)

    add_textbox(
        slide,
        Inches(0.58),
        Inches(6.7),
        Inches(12.0),
        Inches(0.36),
        "Why that audience fit makes sense: the dashboard exposes presets, pantry checkboxes, shopping modes, nearby stores, and plain-language notes. That is the shape of a guided consumer tool, not a specialist analytics console.",
        font_size=14,
        color=MUTED,
    )


def build_slide_3(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide)
    add_top_bar(slide, "Problem", 3)
    add_title(
        slide,
        "The Real-World Problem",
        "Nutrition advice is abstract, but shopping is local, budget-constrained, and full of practical tradeoffs.",
    )

    problems = [
        ("Abstract advice", "Someone may know they want more protein or a lower-cost basket, but that still does not tell them what to buy.", SOFT_BLUE),
        ("Too many tradeoffs", "Calories, protein, produce, convenience, pantry items, and cost all pull the basket in different directions.", SOFT_GREEN),
        ("Local reality", "The decision depends on nearby stores and regional price levels, not just the target on paper.", SOFT_TERRA),
        ("Decision fatigue", "Without structure, the user either guesses or buys a basket that does not really match the goal.", SOFT_GOLD),
    ]
    positions = [(Inches(0.58), Inches(2.15)), (Inches(6.95), Inches(2.15)), (Inches(0.58), Inches(4.55)), (Inches(6.95), Inches(4.55))]
    for (title, body, fill), (left, top) in zip(problems, positions, strict=True):
        add_card(slide, left, top, Inches(5.65), Inches(1.8), fill=fill, line=fill)
        add_textbox(slide, left + Inches(0.2), top + Inches(0.2), Inches(2.4), Inches(0.25), title, font_size=20, bold=True, color=GREEN, font_name=FONT_BOLD)
        add_textbox(slide, left + Inches(0.2), top + Inches(0.62), Inches(5.1), Inches(0.86), body, font_size=16)

    banner = add_card(slide, Inches(0.58), Inches(6.78), Inches(12.1), Inches(0.55), fill=GREEN, line=GREEN)
    banner.text_frame.clear()
    p = banner.text_frame.paragraphs[0]
    p.text = "Communication goal: help a non-expert user turn a health goal into a grocery decision they can actually act on."
    p.font.name = FONT_BOLD
    p.font.size = Pt(17)
    p.font.color.rgb = rgb(CARD)
    style_text_frame(banner, left=12, right=12, top=10, bottom=4, anchor=MSO_ANCHOR.MIDDLE)


def build_slide_4(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide)
    add_top_bar(slide, "How it works", 4)
    add_title(
        slide,
        "What The Tool Does",
        "The clearest explanation is a simple decision pipeline, not a technical architecture diagram.",
    )

    steps = [
        ("1. User inputs", "Location, targets, preferences, pantry, and shopping window."),
        ("2. Basket options", "The planner builds several plausible grocery baskets."),
        ("3. Final choice", "It chooses the basket that best balances nutrition fit, practicality, and price-awareness."),
        ("4. Decision support", "It returns a list, a summary, a rough cost, and store suggestions."),
    ]
    lefts = [Inches(0.62), Inches(3.55), Inches(6.48), Inches(9.41)]
    for idx, (step_title, body) in enumerate(steps):
        left = lefts[idx]
        add_card(slide, left, Inches(2.55), Inches(2.45), Inches(2.35), fill=CARD, line=LINE)
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(0.18), Inches(2.76), Inches(0.48), Inches(0.48))
        circle.fill.solid()
        circle.fill.fore_color.rgb = rgb(SOFT_GOLD)
        circle.line.fill.background()
        tf = circle.text_frame
        style_text_frame(circle, left=0, right=0, top=2, bottom=0, anchor=MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        p.text = str(idx + 1)
        p.alignment = PP_ALIGN.CENTER
        p.font.name = FONT_BOLD
        p.font.size = Pt(18)
        p.font.color.rgb = rgb(GREEN)

        add_textbox(slide, left + Inches(0.16), Inches(3.45), Inches(2.05), Inches(0.3), step_title, font_size=18, bold=True, color=GREEN, font_name=FONT_BOLD)
        add_textbox(slide, left + Inches(0.16), Inches(3.86), Inches(2.05), Inches(0.74), body, font_size=15, color=MUTED)

        if idx < 3:
            arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, left + Inches(2.52), Inches(3.48), Inches(0.42), Inches(0.44))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = rgb(GREEN_2)
            arrow.line.fill.background()

    add_card(slide, Inches(0.62), Inches(5.55), Inches(12.0), Inches(1.02), fill=SOFT_GREEN, line=SOFT_GREEN)
    add_textbox(slide, Inches(0.82), Inches(5.8), Inches(11.6), Inches(0.5), "Plain-English summary", font_size=17, bold=True, color=GREEN, font_name=FONT_BOLD)
    add_textbox(
        slide,
        Inches(0.82),
        Inches(6.14),
        Inches(11.4),
        Inches(0.38),
        "In other words, the app tries a few sensible shopping baskets and picks one that seems closest to the user's goal without ignoring realism.",
        font_size=18,
    )


def build_slide_5(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide)
    add_top_bar(slide, "Interaction", 5)
    add_title(
        slide,
        "How A User Interacts With It",
        "The interface is organized like a guided grocery-planning workflow rather than a technical analytics dashboard.",
    )

    add_image_frame(slide, Inches(0.58), Inches(2.15), Inches(7.55), Inches(3.3), ASSET_DIR / "dashboard-inputs.png")
    add_image_frame(slide, Inches(8.38), Inches(2.15), Inches(4.25), Inches(3.3), ASSET_DIR / "dashboard-stores.png")
    add_textbox(slide, Inches(0.58), Inches(5.52), Inches(2.5), Inches(0.24), "Step 1: set the goal and preferences", font_size=15, bold=True, color=GREEN, font_name=FONT_BOLD)
    add_textbox(slide, Inches(8.38), Inches(5.52), Inches(2.6), Inches(0.24), "Step 2: load nearby supermarkets", font_size=15, bold=True, color=GREEN, font_name=FONT_BOLD)

    add_card(slide, Inches(0.58), Inches(5.9), Inches(12.05), Inches(1.1), fill=CARD, line=LINE)
    add_textbox(slide, Inches(0.78), Inches(6.08), Inches(3.0), Inches(0.24), "Why this is user-friendly", font_size=17, bold=True, color=GREEN, font_name=FONT_BOLD)
    add_bullets(
        slide,
        Inches(0.78),
        Inches(6.35),
        Inches(11.6),
        Inches(0.62),
        [
            "Preset buttons reduce setup effort for common goals like muscle gain, fat loss, and budget-friendly healthy.",
            "The main controls are everyday choices: calories, protein, days, shopping mode, dietary flags, and pantry items.",
            "Nearby store lookup and Recommend are separate steps, which helps a first-time user explore before committing.",
        ],
        font_size=13,
    )


def build_slide_6(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide)
    add_top_bar(slide, "Decision impact", 6)
    add_title(
        slide,
        "What Changes When Decisions Change",
        "These outputs were observed on the deployed app in Mountain View while holding the same base target constant.",
    )

    header = add_card(slide, Inches(0.58), Inches(2.15), Inches(12.05), Inches(0.5), fill=GREEN, line=GREEN)
    style_text_frame(header, left=12, right=12, top=8, bottom=4, anchor=MSO_ANCHOR.MIDDLE)
    p = header.text_frame.paragraphs[0]
    p.text = "Shared setup: Mountain View, CA | 2,200 kcal | 130 g protein | 1 day | balanced shopping mode"
    p.font.name = FONT_BOLD
    p.font.size = Pt(17)
    p.font.color.rgb = rgb(CARD)

    cards = [
        ("Base plan", SCENARIOS["baseline"], SOFT_BLUE),
        ("Budget preference", SCENARIOS["budget_toggle"], SOFT_GREEN),
        ("Vegan preset", SCENARIOS["vegan_preset"], SOFT_TERRA),
    ]
    lefts = [Inches(0.58), Inches(4.65), Inches(8.72)]
    for left, (title, data, fill) in zip(lefts, cards, strict=True):
        add_card(slide, left, Inches(2.95), Inches(3.65), Inches(3.95), fill=fill, line=fill)
        add_textbox(slide, left + Inches(0.2), Inches(3.12), Inches(3.0), Inches(0.24), title, font_size=18, bold=True, color=GREEN, font_name=FONT_BOLD)
        add_textbox(slide, left + Inches(0.2), Inches(3.55), Inches(1.4), Inches(0.42), f"${data['cost']:.2f}", font_size=28, bold=True, font_name=FONT_TITLE)
        add_textbox(slide, left + Inches(0.2), Inches(3.96), Inches(1.8), Inches(0.2), "Typical basket cost", font_size=13, color=MUTED)
        item_text = ", ".join(data["items"][:5])
        add_textbox(slide, left + Inches(0.2), Inches(4.28), Inches(3.15), Inches(0.78), item_text, font_size=13)
        add_textbox(slide, left + Inches(0.2), Inches(5.12), Inches(1.6), Inches(0.2), "Observed effect", font_size=15, bold=True, color=GREEN, font_name=FONT_BOLD)
        add_textbox(slide, left + Inches(0.2), Inches(5.42), Inches(3.12), Inches(0.55), data["note"], font_size=13)
        protein_gap = data["protein_est"] - data["protein_target"]
        add_textbox(slide, left + Inches(0.2), Inches(6.1), Inches(2.8), Inches(0.18), f"Protein gap: {protein_gap:+.1f} g", font_size=13, color=MUTED, bold=True, font_name=FONT_BOLD)
        if "stores" in data:
            add_textbox(slide, left + Inches(0.2), Inches(6.36), Inches(2.8), Inches(0.18), f"Produce fit: {data['stores'].get('produce_pick', 'n/a')}", font_size=13, color=MUTED)

    add_textbox(
        slide,
        Inches(0.58),
        Inches(6.98),
        Inches(12.0),
        Inches(0.24),
        "Takeaway: the app is not just filling calories and protein. It visibly trades off cost, dietary fit, and basket composition when the user changes presets.",
        font_size=13,
        color=MUTED,
    )


def build_slide_7(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide)
    add_top_bar(slide, "Example", 7)
    add_title(
        slide,
        "Key Insight / Recommendation Example",
        "A concrete recommendation helps the audience understand what the tool is actually saying.",
    )

    add_image_frame(slide, Inches(0.58), Inches(2.1), Inches(7.55), Inches(5.1), ASSET_DIR / "dashboard-results.png")
    add_card(slide, Inches(8.38), Inches(2.1), Inches(4.25), Inches(5.1), fill=CARD, line=LINE)
    add_textbox(slide, Inches(8.6), Inches(2.36), Inches(3.6), Inches(0.28), "Observed base recommendation", font_size=18, bold=True, color=GREEN, font_name=FONT_BOLD)
    add_textbox(slide, Inches(8.6), Inches(2.78), Inches(2.4), Inches(0.22), "What the model is saying", font_size=15, color=MUTED, bold=True, font_name=FONT_BOLD)
    add_textbox(
        slide,
        Inches(8.6),
        Inches(3.12),
        Inches(3.55),
        Inches(1.15),
        "Buy a small set of versatile staples that gets you close to the target without making the trip expensive or overly complicated.",
        font_size=19,
    )
    add_labeled_metric(slide, Inches(8.6), Inches(4.3), Inches(3.55), Inches(0.46), "Protein", f"{SCENARIOS['baseline']['protein_est']:.1f} g vs {SCENARIOS['baseline']['protein_target']:.1f} g target")
    add_labeled_metric(slide, Inches(8.6), Inches(4.92), Inches(3.55), Inches(0.46), "Calories", f"{SCENARIOS['baseline']['calorie_est']:.1f} vs {SCENARIOS['baseline']['calorie_target']:.1f} target")
    add_labeled_metric(slide, Inches(8.6), Inches(5.54), Inches(3.55), Inches(0.46), "Cost", f"About ${SCENARIOS['baseline']['cost']:.2f} for the basket")
    add_textbox(slide, Inches(8.6), Inches(6.18), Inches(2.0), Inches(0.22), "Why this matters", font_size=15, color=MUTED, bold=True, font_name=FONT_BOLD)
    add_textbox(
        slide,
        Inches(8.6),
        Inches(6.48),
        Inches(3.55),
        Inches(0.66),
        "The recommendation is useful because it is actionable: it tells the user what to buy, what the basket roughly achieves, and where it may fit best nearby. In this observed case, Safeway is the one-stop pick and La Plaza is the stronger produce fit.",
        font_size=13,
    )


def build_slide_8(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide)
    add_top_bar(slide, "Limits", 8)
    add_title(
        slide,
        "Limitations + How To Improve",
        "Strong communication includes being honest about uncertainty, assumptions, and what the current tool still cannot do.",
    )

    add_card(slide, Inches(0.58), Inches(2.15), Inches(5.55), Inches(4.0), fill=SOFT_TERRA, line=SOFT_TERRA)
    add_textbox(slide, Inches(0.78), Inches(2.42), Inches(2.8), Inches(0.24), "Current limitations", font_size=18, bold=True, color=GREEN, font_name=FONT_BOLD)
    add_bullets(
        slide,
        Inches(0.78),
        Inches(2.72),
        Inches(5.0),
        Inches(3.15),
        [
            "The basket uses generic foods, not exact branded products or exact live store inventory.",
            "Price estimates are regional references from public price tables, not store quotes.",
            "Store recommendations are coarse fit suggestions based on nearby store type, brand, and distance.",
            "Long shopping windows are heuristic. Fresh and bulk modes both soften perishables, so results become approximate.",
            "The tool is not medical nutrition advice and does not capture allergies or complex health needs.",
        ],
        font_size=14,
    )

    add_card(slide, Inches(6.48), Inches(2.15), Inches(6.15), Inches(4.0), fill=SOFT_GREEN, line=SOFT_GREEN)
    add_textbox(slide, Inches(6.68), Inches(2.42), Inches(2.8), Inches(0.24), "How to improve it", font_size=18, bold=True, color=GREEN, font_name=FONT_BOLD)
    add_bullets(
        slide,
        Inches(6.68),
        Inches(2.72),
        Inches(5.6),
        Inches(3.15),
        [
            "Add real store inventory and current store-specific prices where possible.",
            "Let the user set a hard spending cap, not just a budget-friendly preference.",
            "Support allergies, household size, cultural preferences, and stronger nutrition guardrails.",
            "Collect feedback from real student or young professional users and adapt the interface around what confuses them.",
            "Show confidence labels more clearly so users know what is exact versus approximate.",
        ],
        font_size=14,
    )

    evidence = (
        f"Observed evidence: in 7-day mode, fresh landed at {SCENARIOS['fresh7']['protein_est']:.1f} g protein versus a "
        f"{SCENARIOS['fresh7']['protein_target']:.1f} g target, and bulk landed at {SCENARIOS['bulk7']['protein_est']:.1f} g "
        f"versus {SCENARIOS['bulk7']['protein_target']:.1f} g. That shows why perishables and scaling remain a real limitation."
    )
    add_card(slide, Inches(0.58), Inches(6.4), Inches(12.05), Inches(0.58), fill=CARD, line=LINE)
    add_textbox(slide, Inches(0.78), Inches(6.56), Inches(11.7), Inches(0.24), evidence, font_size=13, color=MUTED)


def build_slide_9(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide)
    add_top_bar(slide, "Ethics", 9)
    add_title(
        slide,
        "Stakeholders + Ethics",
        "A useful recommendation system should help without overstating certainty or excluding the people who need help most.",
    )

    stakeholder_boxes = [
        (Inches(0.58), Inches(2.2), "Shoppers", "They gain a clearer starting point for cost, nutrition, and store choice, especially if they feel overwhelmed by grocery planning.", SOFT_BLUE),
        (Inches(7.0), Inches(2.2), "Stores and educators", "Nearby stores may benefit from more intentional trips, and nutrition coaches or campus wellness groups could use the tool as a teaching aid.", SOFT_GREEN),
        (Inches(0.58), Inches(4.1), "Potentially underserved users", "People with allergies, medical conditions, low food access, or strong cultural constraints may need more safeguards than the current tool provides.", SOFT_TERRA),
        (Inches(7.0), Inches(4.1), "Ethical design responsibility", "The tool should keep telling the truth about approximation: regional prices are not store quotes and store-fit suggestions are not inventory guarantees.", SOFT_GOLD),
    ]
    for left, top, title, body, fill in stakeholder_boxes:
        add_card(slide, left, top, Inches(5.8), Inches(1.45), fill=fill, line=fill)
        add_textbox(slide, left + Inches(0.2), top + Inches(0.2), Inches(3.3), Inches(0.24), title, font_size=18, bold=True, color=GREEN, font_name=FONT_BOLD)
        add_textbox(slide, left + Inches(0.2), top + Inches(0.62), Inches(5.2), Inches(0.56), body, font_size=15)

    footer = add_card(slide, Inches(0.58), Inches(6.15), Inches(12.05), Inches(0.88), fill=GREEN, line=GREEN)
    style_text_frame(footer, left=12, right=12, top=8, bottom=4)
    footer.text_frame.clear()
    p1 = footer.text_frame.paragraphs[0]
    p1.text = "Thoughtful ethical stance"
    p1.font.name = FONT_BOLD
    p1.font.size = Pt(17)
    p1.font.color.rgb = rgb(CARD)
    p2 = footer.text_frame.add_paragraph()
    p2.text = "The responsible position is to present the tool as a practical starting point for grocery decisions, not as medical advice or a perfect answer. Transparency is part of the product, not a disclaimer added at the end."
    p2.font.name = FONT_BODY
    p2.font.size = Pt(14)
    p2.font.color.rgb = rgb(CARD)


def build_slide_10(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide)
    add_top_bar(slide, "Close", 10)
    add_title(
        slide,
        "Takeaways + Future Use",
        "The strongest close is simple: this project makes nutrition goals more actionable for an everyday shopper.",
    )

    takeaways = [
        (Inches(0.58), "Why it matters", "It turns an abstract nutrition goal into groceries, prices, and store choices a person can immediately use.", SOFT_BLUE),
        (Inches(4.63), "Why it is worth using", "The outputs are understandable, local-feeling, and honest about what is approximate.", SOFT_GREEN),
        (Inches(8.68), "What comes next", "Pilot it with real users, collect feedback, and strengthen the data and confidence signals.", SOFT_TERRA),
    ]
    for left, title, body, fill in takeaways:
        add_card(slide, left, Inches(2.55), Inches(3.75), Inches(2.7), fill=fill, line=fill)
        add_textbox(slide, left + Inches(0.22), Inches(2.86), Inches(3.0), Inches(0.3), title, font_size=19, bold=True, color=GREEN, font_name=FONT_BOLD)
        add_textbox(slide, left + Inches(0.22), Inches(3.32), Inches(3.2), Inches(1.45), body, font_size=17)

    footer = add_card(slide, Inches(0.58), Inches(5.8), Inches(12.05), Inches(0.95), fill=GREEN, line=GREEN)
    style_text_frame(footer, left=12, right=12, top=8, bottom=4)
    footer.text_frame.clear()
    p1 = footer.text_frame.paragraphs[0]
    p1.text = "Closing line"
    p1.font.name = FONT_BOLD
    p1.font.size = Pt(17)
    p1.font.color.rgb = rgb(CARD)
    p2 = footer.text_frame.add_paragraph()
    p2.text = "If the goal is better grocery decisions for non-experts, this tool already works as a strong communication bridge between nutrition targets and real shopping behavior."
    p2.font.name = FONT_BODY
    p2.font.size = Pt(16)
    p2.font.color.rgb = rgb(CARD)


def build_slide_11(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_background(slide)
    add_top_bar(slide, "Appendix", 11)
    add_title(
        slide,
        "Demo Backup",
        "If the live demo fails, this backup slide still lets the audience see the full recommendation workflow and outputs.",
    )

    add_image_frame(slide, Inches(0.58), Inches(2.05), Inches(8.85), Inches(5.25), ASSET_DIR / "dashboard-full-results.png")
    add_card(slide, Inches(9.7), Inches(2.05), Inches(2.93), Inches(5.25), fill=CARD, line=LINE)
    add_textbox(slide, Inches(9.9), Inches(2.32), Inches(2.4), Inches(0.28), "What I would point to", font_size=17, bold=True, color=GREEN, font_name=FONT_BOLD)
    add_bullets(
        slide,
        Inches(9.9),
        Inches(2.72),
        Inches(2.4),
        Inches(3.9),
        [
            "The grouped shopping list by food role.",
            "The cost estimate and price notes.",
            "The recommended store picks for the basket.",
            "The nutrition summary versus target.",
            "The approximation notes that prevent overconfidence.",
        ],
        font_size=14,
    )


def add_notes_from_transcript(prs: Presentation) -> None:
    scripts = [
        "Today I am presenting a grocery-planning dashboard that takes a nutrition goal and turns it into a practical shopping trip.",
        "The clearest audience for this tool is a busy, health-conscious shopper who is not a nutrition expert.",
        "The real-world problem is that eating for a goal sounds simple until someone has to stand in a store and decide what to buy.",
        "At a high level, the tool asks for a location, a daily target, a few food preferences, and any pantry items the user already has.",
        "This slide shows why the tool is accessible for a non-specialist audience.",
        "Here is the most important interaction story: changing one preference changes cost, basket composition, and sometimes store fit.",
        "This is a concrete example of the tool speaking in plain English through a single observed recommendation.",
        "The tool is useful, but it is important to state its limits clearly and honestly.",
        "Several groups are affected by this tool, so transparency and accessibility matter.",
        "The big takeaway is that this project makes nutrition goals more actionable.",
        "This backup slide is here in case a live demo fails.",
    ]
    for slide, script in zip(prs.slides, scripts, strict=True):
        try:
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
        except Exception:
            continue
        notes_text_frame.text = script


def build_editable_presentation() -> Path:
    OUT_DIR.mkdir(exist_ok=True)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    prs.core_properties.title = "From Nutrition Goals to a Grocery Trip"
    prs.core_properties.subject = "Editable INDENG 243 final presentation"
    prs.core_properties.author = "Codex"

    build_slide_1(prs)
    build_slide_2(prs)
    build_slide_3(prs)
    build_slide_4(prs)
    build_slide_5(prs)
    build_slide_6(prs)
    build_slide_7(prs)
    build_slide_8(prs)
    build_slide_9(prs)
    build_slide_10(prs)
    build_slide_11(prs)
    add_notes_from_transcript(prs)

    prs.save(OUT_FILE)
    return OUT_FILE


def main() -> None:
    path = build_editable_presentation()
    print(path)


if __name__ == "__main__":
    main()
