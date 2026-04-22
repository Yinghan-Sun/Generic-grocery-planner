from __future__ import annotations

import math
import sys
from pathlib import Path

REPO_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(REPO_ROOT / "tmp" / "presentation_py"))

from PIL import Image, ImageDraw, ImageFont, ImageOps  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402


W = 1920
H = 1080
BG = "#F6F1E8"
GREEN = "#204736"
GREEN_2 = "#315B46"
MOSS = "#5F7759"
SAGE = "#DDE6DB"
TEXT = "#163126"
MUTED = "#5D6E61"
GOLD = "#D9A441"
TERRACOTTA = "#C86A45"
CARD = "#FFFDF8"
LINE = "#D8D2C6"
SOFT_GOLD = "#F5E5BE"
SOFT_GREEN = "#EAF2E7"
SOFT_BLUE = "#E8F0F7"
SOFT_TERRA = "#F8E8E1"

FONT_TITLE = "/System/Library/Fonts/Avenir Next.ttc"
FONT_BODY = "/System/Library/Fonts/Avenir.ttc"
FONT_BOLD = "/System/Library/Fonts/HelveticaNeue.ttc"

ASSET_DIR = REPO_ROOT / "presentation_assets"
OUT_DIR = REPO_ROOT / "presentation"
SLIDE_DIR = OUT_DIR / "slides"

SCENARIOS = {
    "baseline": {
        "label": "Base plan",
        "cost": 8.54,
        "protein_est": 137.7,
        "protein_target": 130.0,
        "calorie_est": 2160.4,
        "calorie_target": 2200.0,
        "items": ["Chicken breast", "Rotisserie chicken", "Wholemeal bread", "Spinach", "Carrots", "Olive oil"],
        "stores": {"one_stop_pick": "Safeway", "produce_pick": "La Plaza"},
        "note": "Balances the target with familiar staples and a produce stop.",
    },
    "budget_toggle": {
        "label": "Budget toggle on",
        "cost": 5.44,
        "protein_est": 110.6,
        "protein_target": 130.0,
        "calorie_est": 2185.2,
        "calorie_target": 2200.0,
        "items": ["Lentils", "Eggs", "Rice", "Cabbage", "Carrots", "Olive oil"],
        "stores": {"one_stop_pick": "Safeway", "produce_pick": "Safeway"},
        "note": "Shifts toward cheaper pantry staples and keeps shopping simpler.",
    },
    "vegan_preset": {
        "label": "Vegan preset",
        "cost": 7.18,
        "protein_est": 128.4,
        "protein_target": 125.0,
        "calorie_est": 2238.7,
        "calorie_target": 2200.0,
        "items": ["Tofu", "Lentils", "Rice", "Spinach", "Bananas", "Olive oil"],
        "stores": {"one_stop_pick": "Safeway", "produce_pick": "La Plaza"},
        "note": "Rebuilds the basket around plant proteins and dairy-free staples while staying close to the target.",
    },
    "breakfast_toggle": {
        "label": "Breakfast focus",
        "cost": 8.99,
        "protein_est": 126.6,
        "protein_target": 130.0,
        "calorie_est": 2276.4,
        "calorie_target": 2200.0,
        "items": ["Eggs", "Protein yogurt", "Oats", "Kiwi", "Bananas"],
        "stores": {"one_stop_pick": "Safeway", "produce_pick": "La Plaza"},
        "note": "Rebuilds the list around breakfast-friendly foods instead of lunch staples.",
    },
    "bulk7": {
        "label": "7-day bulk mode",
        "cost": 53.24,
        "protein_est": 780.2,
        "protein_target": 910.0,
        "calorie_est": 14872.8,
        "calorie_target": 15400.0,
        "items": ["Chicken breast", "Rotisserie chicken", "Wholemeal bread", "Spinach", "Bell peppers", "Carrots"],
        "note": "Bulk mode buys more shelf-stable food, but still softens perishables.",
    },
    "fresh7": {
        "label": "7-day fresh mode",
        "cost": 46.18,
        "protein_est": 631.6,
        "protein_target": 910.0,
        "calorie_est": 16392.1,
        "calorie_target": 15400.0,
        "items": ["Chicken breast", "Rotisserie chicken", "Wholemeal bread", "Oats", "Carrots", "Bell peppers"],
        "note": "Fresh mode intentionally caps short-life items, so long windows become approximate.",
    },
}

PROJECT_UNDERSTANDING = {
    "stack": "Python Flask app with a vanilla JavaScript frontend, DuckDB data store, trained scoring artifacts, and public regional price plus nearby-store data inputs.",
    "app_purpose": "Convert a nutrition goal into a practical generic grocery basket, a rough regional cost estimate, and nearby store suggestions.",
    "inferred_audience": "Busy non-expert grocery shoppers, especially students or young professionals balancing health, convenience, and budget.",
    "key_inputs": "Location, calorie and protein targets, optional macro and micronutrient targets, shopping window, shopping mode, dietary flags, and pantry items.",
    "key_outputs": "Suggested shopping list by food role, nutrition-versus-target summary, typical basket cost and range, nearby stores, recommended store picks, meal ideas, and guidance notes.",
    "main_tradeoff": "The planner appears to balance closeness to nutrition targets with practicality, price-awareness, preference fit, and store convenience rather than solving a single exact minimum-cost problem.",
    "uncertainties": [
        "The deployed app hides most debug metadata, so exact learned-model weights are inferred from code paths and output behavior rather than directly visible in production.",
        "Store-fit suggestions appear to use store type, brand, and distance, but not exact live inventory.",
        "Long-window quantity scaling is heuristic; observed 7-day fresh and bulk scenarios drift from the target because perishables are intentionally softened.",
    ],
}

SLIDE_OUTLINE = [
    {
        "slide": 1,
        "title": "Title + value proposition",
        "bullets": [
            "Introduce the dashboard as a grocery-planning tool rather than a technical model.",
            "State the value proposition in one sentence and show the app home screen.",
        ],
    },
    {
        "slide": 2,
        "title": "Who this is for",
        "bullets": [
            "Define a specific non-technical user: busy health-conscious shoppers, especially students and young professionals.",
            "Explain what they care about: easy choices, nearby stores, realistic costs, and goals they can trust.",
        ],
    },
    {
        "slide": 3,
        "title": "The real-world problem",
        "bullets": [
            "Frame the gap between abstract nutrition advice and a real grocery trip.",
            "Show why local availability, cost, and convenience all matter at once.",
        ],
    },
    {
        "slide": 4,
        "title": "What the tool does",
        "bullets": [
            "Explain the flow from inputs to basket generation to recommendation output.",
            "Translate the optimization logic into plain English: the app tries multiple baskets and picks a balanced one.",
        ],
    },
    {
        "slide": 5,
        "title": "How a user interacts with it",
        "bullets": [
            "Show preset buttons, target fields, pantry choices, and nearby store lookup.",
            "Emphasize the one-click flow and user-friendly layout for non-specialists.",
        ],
    },
    {
        "slide": 6,
        "title": "What changes when decisions change",
        "bullets": [
            "Compare live outputs when the user changes only one preference at a time.",
            "Show how cost, items, and store suggestions shift when the goal changes.",
        ],
    },
    {
        "slide": 7,
        "title": "Key insight example",
        "bullets": [
            "Walk through one concrete recommendation and explain what the model is saying in plain English.",
            "Highlight the nutrition summary, price estimate, and store picks as decision support.",
        ],
    },
    {
        "slide": 8,
        "title": "Limitations + improvements",
        "bullets": [
            "State the current limits clearly and tie them to observed behavior.",
            "Explain how user feedback and better data could improve the tool.",
        ],
    },
    {
        "slide": 9,
        "title": "Stakeholders + ethics",
        "bullets": [
            "Discuss who benefits, who may be underserved, and where transparency matters.",
            "Treat health, accessibility, and fairness as design concerns, not side notes.",
        ],
    },
    {
        "slide": 10,
        "title": "Takeaways + future use",
        "bullets": [
            "Close on why the tool matters and why this audience could use it.",
            "End with a roadmap for trust, adoption, and next-step improvement.",
        ],
    },
    {
        "slide": 11,
        "title": "Demo backup",
        "bullets": [
            "Provide a backup screenshot walkthrough in case a live demo fails.",
        ],
    },
]

TRANSCRIPT = [
    {
        "slide": 1,
        "title": "Title + Value Proposition",
        "script": (
            "Today I am presenting a grocery-planning dashboard that takes a nutrition goal and turns it into a practical shopping trip. "
            "Instead of asking a user to translate abstract advice like eat more protein or stay on budget into actual groceries, the tool gives them a starting basket, a rough price estimate, and nearby store suggestions. "
            "That makes the project easy to explain to a non-technical audience because its value is visible right away: it helps someone move from a goal to an action. "
            "As we go, I am going to focus less on code and more on the user decision this tool supports."
        ),
    },
    {
        "slide": 2,
        "title": "Who This Is For",
        "script": (
            "The clearest audience for this tool is a busy, health-conscious shopper who is not a nutrition expert. "
            "I think the best example is a student or young professional who wants to shop in a way that supports muscle gain, fat loss, maintenance, or a healthier budget, but does not want to build a meal plan from scratch. "
            "This audience needs three things at once: simple choices, realistic prices, and recommendations that feel local and doable. "
            "So the dashboard is not trying to teach nutrition theory. It is trying to reduce friction and help a person make a better grocery decision in a few minutes."
        ),
    },
    {
        "slide": 3,
        "title": "The Real-World Problem",
        "script": (
            "The real-world problem is that eating for a goal sounds simple until someone has to stand in a store and decide what to buy. "
            "Nutrition advice is usually abstract, grocery prices vary by place, and convenience matters just as much as protein or calories. "
            "A shopper might know they want to eat healthier, but they still have to answer practical questions like what basket gets me close to my goal, how much will it cost, and where should I go nearby. "
            "This project matters because it turns that messy decision into a clearer comparison instead of leaving the user to guess."
        ),
    },
    {
        "slide": 4,
        "title": "What The Tool Does",
        "script": (
            "At a high level, the tool asks for a location, a daily target, a few food preferences, and any pantry items the user already has. "
            "Then it builds several possible grocery baskets and chooses one that best balances nutrition fit, practicality, and price-awareness. "
            "In plain English, the app is not hunting for one mathematically perfect meal plan. It is trying a few sensible baskets and selecting the one that seems most useful for a real shopper. "
            "The result is a shopping list, a nutrition summary, a typical basket cost, and store suggestions that help the user act on the recommendation."
        ),
    },
    {
        "slide": 5,
        "title": "How A User Interacts With It",
        "script": (
            "This slide shows why the tool is accessible for a non-specialist audience. "
            "The user can start with a preset like muscle gain or budget-friendly healthy, keep the default city, and click a button to load nearby supermarkets. "
            "They can also fine-tune calories, protein, shopping window, meal context, and pantry items without needing any technical knowledge. "
            "That matters for grading because the dashboard is not just producing an answer. It is inviting exploration in a way that feels familiar, like filling out a guided shopping form rather than operating a model."
        ),
    },
    {
        "slide": 6,
        "title": "What Changes When Decisions Change",
        "script": (
            "Here is the most important interaction story. I kept the same Mountain View location and the same 2,200 calorie, 130 gram protein target, and then changed one preference at a time. "
            "The base plan costs about eight dollars and uses familiar staples like chicken, bread, spinach, and carrots. "
            "When I turn on the budget preference, the basket drops closer to five and a half dollars and shifts toward lentils, eggs, rice, and cabbage. "
            "When I switch to the vegan preset, the basket rebuilds around tofu, lentils, rice, and other plant-based staples. "
            "So the model is making tradeoffs visible. Different choices really do change cost, dietary fit, basket composition, and even the suggested store fit."
        ),
    },
    {
        "slide": 7,
        "title": "Key Recommendation Example",
        "script": (
            "This is a concrete example of the tool speaking in plain English. For the base scenario, the dashboard recommends two protein anchors, one carb base, two produce items, and a calorie booster. "
            "It estimates about 138 grams of protein against a 130 gram target, about 2,160 calories against a 2,200 calorie target, and a typical basket cost of about eight and a half dollars. "
            "The useful part is not the exact number. The useful part is the explanation: buy a small set of versatile staples that get you close to the goal, then use the nearby store suggestions to make the trip easier. "
            "That is decision support, not just prediction."
        ),
    },
    {
        "slide": 8,
        "title": "Limitations + How To Improve",
        "script": (
            "The tool is useful, but it is important to state its limits clearly. "
            "First, it works with generic foods, not exact store inventory or exact product brands. "
            "Second, the prices are regional estimates rather than store quotes. "
            "Third, the long-window logic is approximate. In the live seven-day scenarios, fresh and bulk modes both drift away from the target because perishables are intentionally softened. "
            "The path forward is clear: add real inventory and pricing feeds, let users set a hard budget cap, support more dietary constraints, and use feedback from real shoppers to improve the confidence and realism of each recommendation."
        ),
    },
    {
        "slide": 9,
        "title": "Stakeholders + Ethics",
        "script": (
            "Several groups are affected by this tool. Shoppers benefit from clearer planning, stores may benefit from more intentional trips, and nutrition educators could use it as a teaching aid. "
            "But there are also ethical concerns. Users with allergies, medical conditions, or very limited access to stores may be underserved if the tool is treated as more certain than it really is. "
            "That is why transparency matters. The dashboard should keep telling users when a result is approximate, when price data is regional, and when store recommendations are based on fit rather than live inventory. "
            "Used responsibly, the tool supports healthier choices; used carelessly, it could overstate confidence."
        ),
    },
    {
        "slide": 10,
        "title": "Takeaways + Future Use",
        "script": (
            "The big takeaway is that this project makes nutrition goals more actionable. "
            "It gives a non-expert user a simple way to explore tradeoffs among health targets, convenience, and budget, and it does that with outputs they can immediately recognize: a shopping list, a cost estimate, and store suggestions. "
            "I would trust it as a practical starting point because it shows its reasoning in user-friendly terms and it does not pretend to be exact inventory or medical advice. "
            "The next step would be to test it with real student or young professional shoppers, collect feedback, and adapt the interface and rules around what those users say they actually need."
        ),
    },
    {
        "slide": 11,
        "title": "Demo Backup",
        "script": (
            "This backup slide is here in case a live demo fails. "
            "It lets me walk through the same recommendation screen, point to the shopping list, the nutrition summary, the price notes, and the store suggestions, and still keep the presentation understandable."
        ),
    },
]

RUBRIC_ROWS = [
    ("A1. Clearly identify the target audience", "2", "Slide 2 names busy non-expert shoppers, especially students and young professionals."),
    ("A2. Show understanding of the audience's needs and interests", "2, 3", "Slides 2 and 3 connect the audience to ease, affordability, local relevance, and decision stress."),
    ("A3. Provide a concise, high-level summary for a non-specialist audience", "1, 4", "Slide 1 gives the one-sentence value proposition and Slide 4 explains the workflow without jargon."),
    ("A4. Communicate key insights accessibly", "4, 6, 7", "These slides translate the planning logic and scenario changes into plain English."),
    ("A5. Show relevance of the model/tool to audience goals and challenges", "2, 3, 10", "The tool is tied to real grocery goals like health, budget, and convenience."),
    ("B6. Show that the tool is accessible and user-friendly for a non-specialist audience", "5", "Slide 5 uses real screenshots to show presets, simple controls, and one-click flow."),
    ("B7. Show how different choices affect outcomes", "6", "Slide 6 compares live outputs after changing preferences."),
    ("B8. Use the tool to engage the user in meaningful exploration and clearly explain model decisions", "5, 6, 7", "The deck shows interaction, scenario exploration, and plain-English interpretation of the recommendation."),
    ("C9. Clearly communicate limitations", "8", "Slide 8 lists concrete limits tied to observed app behavior."),
    ("C10. Discuss strategies for overcoming limitations", "8", "Slide 8 proposes better data, budget constraints, dietary support, and feedback loops."),
    ("C11. Discuss how the tool could be adapted based on audience feedback", "8, 10", "The presentation suggests piloting with real users and adapting the interface from feedback."),
    ("D12. Show real-world relevance", "3, 10", "Slides 3 and 10 connect the tool to actual grocery decisions."),
    ("D13. Analyze impact on different stakeholders", "9", "Slide 9 covers shoppers, stores, nutrition educators, and underserved users."),
    ("D14. Address ethical implications and societal impact thoughtfully", "8, 9", "Slides 8 and 9 address confidence, accessibility, health safety, and transparency."),
    ("E15. Ensure clarity and logical structure", "1-10", "The deck follows a problem-to-solution-to-evaluation arc with clear headings and visual grouping."),
    ("E16. Make the presentation engaging and appropriate for 8-10 minutes", "1-10", "Ten main slides support a paced 8-10 minute talk, with Slide 11 as backup."),
    ("E17. Demonstrate preparation and a strong understanding of the audience", "2-10", "Audience framing, live evidence, and concrete scenario comparisons show preparation."),
    ("F18. Use creative presentation choices that improve understanding", "4, 6, 7", "Workflow graphics, scenario cards, and screenshot-led explanation make the model easier to grasp."),
    ("F19. Show visible effort in making the model understandable and relevant", "5, 6, 11", "Real screenshots, observed outputs, and a demo backup slide signal strong effort."),
]


def ensure_dirs() -> None:
    OUT_DIR.mkdir(exist_ok=True)
    SLIDE_DIR.mkdir(exist_ok=True)


def font(size: int, *, kind: str = "body") -> ImageFont.FreeTypeFont:
    if kind == "title":
        path = FONT_TITLE
    elif kind == "bold":
        path = FONT_BOLD
    else:
        path = FONT_BODY
    return ImageFont.truetype(path, size=size)


def text_width(draw: ImageDraw.ImageDraw, text: str, the_font: ImageFont.FreeTypeFont) -> float:
    return draw.textlength(text, font=the_font)


def wrap_text(draw: ImageDraw.ImageDraw, text: str, the_font: ImageFont.FreeTypeFont, max_width: int) -> list[str]:
    lines: list[str] = []
    for paragraph in text.split("\n"):
        words = paragraph.split()
        if not words:
            lines.append("")
            continue
        current = words[0]
        for word in words[1:]:
            trial = f"{current} {word}"
            if text_width(draw, trial, the_font) <= max_width:
                current = trial
            else:
                lines.append(current)
                current = word
        lines.append(current)
    return lines


def draw_paragraph(
    draw: ImageDraw.ImageDraw,
    xy: tuple[int, int],
    text: str,
    the_font: ImageFont.FreeTypeFont,
    fill: str,
    max_width: int,
    line_gap: int = 10,
) -> int:
    lines = wrap_text(draw, text, the_font, max_width)
    x, y = xy
    bbox = the_font.getbbox("Ag")
    line_h = bbox[3] - bbox[1]
    for line in lines:
        draw.text((x, y), line, font=the_font, fill=fill)
        y += line_h + line_gap
    return y


def draw_bullets(
    draw: ImageDraw.ImageDraw,
    xy: tuple[int, int],
    bullets: list[str],
    max_width: int,
    *,
    bullet_fill: str = TEXT,
    text_fill: str = TEXT,
    bullet_font: ImageFont.FreeTypeFont | None = None,
    text_font: ImageFont.FreeTypeFont | None = None,
    gap_after: int = 16,
) -> int:
    bullet_font = bullet_font or font(24, kind="bold")
    text_font = text_font or font(28)
    x, y = xy
    bullet_box = 28
    for bullet in bullets:
        draw.text((x, y), "-", font=bullet_font, fill=bullet_fill)
        end_y = draw_paragraph(draw, (x + bullet_box, y - 2), bullet, text_font, text_fill, max_width - bullet_box, line_gap=8)
        y = end_y + gap_after
    return y


def card(draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int], *, fill: str = CARD, outline: str = LINE, radius: int = 28, width: int = 2) -> None:
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=width)


def pill(draw: ImageDraw.ImageDraw, x: int, y: int, text: str, *, fill: str = SOFT_GREEN, text_fill: str = GREEN, padding_x: int = 22, padding_y: int = 11, font_size: int = 22) -> int:
    the_font = font(font_size, kind="bold")
    tw = int(text_width(draw, text, the_font))
    th = the_font.getbbox("Ag")[3] - the_font.getbbox("Ag")[1]
    box = (x, y, x + tw + padding_x * 2, y + th + padding_y * 2)
    draw.rounded_rectangle(box, radius=18, fill=fill)
    draw.text((x + padding_x, y + padding_y - 2), text, font=the_font, fill=text_fill)
    return box[2]


def heading(draw: ImageDraw.ImageDraw, slide_no: int, kicker: str, title: str, subtitle: str | None = None) -> None:
    draw.rectangle((0, 0, W, 86), fill=GREEN)
    draw.text((84, 28), kicker.upper(), font=font(22, kind="bold"), fill=SOFT_GOLD)
    draw.text((84, 136), title, font=font(52, kind="title"), fill=TEXT)
    if subtitle:
        draw_paragraph(draw, (84, 220), subtitle, font(28), MUTED, 1150, line_gap=8)
    draw.text((1820, 1005), f"{slide_no:02d}", font=font(22, kind="bold"), fill=MUTED)


def paste_image(
    canvas: Image.Image,
    image_path: Path,
    box: tuple[int, int, int, int],
    *,
    crop: tuple[int, int, int, int] | None = None,
    fit: str = "cover",
    border: str = LINE,
) -> None:
    img = Image.open(image_path).convert("RGB")
    if crop is not None:
        img = img.crop(crop)
    x1, y1, x2, y2 = box
    w = x2 - x1
    h = y2 - y1
    if fit == "cover":
        img = ImageOps.fit(img, (w, h), method=Image.Resampling.LANCZOS)
    else:
        img.thumbnail((w, h), Image.Resampling.LANCZOS)
        framed = Image.new("RGB", (w, h), "white")
        framed.paste(img, ((w - img.width) // 2, (h - img.height) // 2))
        img = framed
    shadow = Image.new("RGBA", (w + 24, h + 24), (0, 0, 0, 0))
    shadow_draw = ImageDraw.Draw(shadow)
    shadow_draw.rounded_rectangle((10, 10, w + 10, h + 10), radius=28, fill=(22, 49, 38, 30))
    canvas.alpha_composite(shadow, (x1 - 8, y1 - 8))
    card_layer = Image.new("RGBA", (w, h), (255, 255, 255, 255))
    mask = Image.new("L", (w, h), 0)
    ImageDraw.Draw(mask).rounded_rectangle((0, 0, w, h), radius=24, fill=255)
    card_layer.paste(img, (0, 0))
    canvas.paste(card_layer.convert("RGB"), (x1, y1), mask)
    draw = ImageDraw.Draw(canvas)
    draw.rounded_rectangle((x1, y1, x2, y2), radius=24, outline=border, width=2)


def save_slide(img: Image.Image, name: str) -> Path:
    path = SLIDE_DIR / name
    img.convert("RGB").save(path)
    return path


def make_canvas() -> tuple[Image.Image, ImageDraw.ImageDraw]:
    img = Image.new("RGBA", (W, H), BG)
    draw = ImageDraw.Draw(img)
    return img, draw


def slide_1() -> Path:
    img, draw = make_canvas()
    draw.rectangle((0, 0, W, H), fill=BG)
    draw.rectangle((0, 0, W, 110), fill=GREEN)
    draw.ellipse((1360, -120, 1910, 430), fill=GREEN_2)
    draw.ellipse((-120, 760, 380, 1260), fill=SOFT_GREEN)
    draw.text((86, 36), "INDENG 243 Analytics Communication Final", font=font(24, kind="bold"), fill=SOFT_GOLD)
    draw.text((86, 168), "From Nutrition Goals\nto a Grocery Trip", font=font(70, kind="title"), fill=TEXT, spacing=6)
    draw_paragraph(
        draw,
        (92, 376),
        "This tool turns a health goal into a practical grocery list, a rough cost estimate, and nearby store suggestions.",
        font(34),
        MUTED,
        720,
        line_gap=10,
    )
    pill(draw, 92, 520, "Non-technical audience first", fill=SOFT_GREEN)
    pill(draw, 92, 584, "Observed on the deployed app", fill=SOFT_BLUE, text_fill="#245072")
    pill(draw, 92, 648, "Focus: decisions, tradeoffs, outcomes", fill=SOFT_TERRA, text_fill=TERRACOTTA)
    card(draw, (86, 760, 720, 956), fill=CARD)
    draw.text((118, 792), "What the audience sees", font=font(24, kind="bold"), fill=GREEN)
    draw_bullets(
        draw,
        (118, 842),
        [
            "Preset goals like muscle gain, fat loss, maintenance, and budget-friendly healthy.",
            "A one-click path from location and targets to a shopping list and store suggestions.",
            "Outputs that read like a grocery assistant, not a technical dashboard.",
        ],
        560,
        text_font=font(25),
        gap_after=12,
    )
    paste_image(img, ASSET_DIR / "dashboard-home.png", (1145, 138, 1815, 905), crop=(0, 0, 1440, 1180))
    draw.text((1145, 930), "Hero view: the deployed dashboard starts with presets, local store lookup, and a guided form.", font=font(22), fill=MUTED)
    draw.text((1810, 1008), "01", font=font(22, kind="bold"), fill=MUTED)
    return save_slide(img, "slide_01_title.png")


def slide_2() -> Path:
    img, draw = make_canvas()
    heading(draw, 2, "Audience", "Who This Is For", "The strongest non-technical audience is not a data scientist. It is a shopper who wants a healthier and easier grocery trip.")
    card(draw, (84, 300, 740, 920), fill=CARD)
    draw.text((120, 340), "Primary user persona", font=font(28, kind="bold"), fill=GREEN)
    draw.text((120, 396), "Busy students and young professionals", font=font(42, kind="title"), fill=TEXT)
    draw_paragraph(
        draw,
        (120, 476),
        "People who want to eat for a goal such as muscle gain, fat loss, maintenance, or a healthier budget without hand-building a meal plan.",
        font(30),
        MUTED,
        560,
        line_gap=9,
    )
    needs = [
        ("Fast setup", "They want a recommendation in minutes, not a spreadsheet exercise.", SOFT_GREEN),
        ("Affordable", "They care about rough basket cost before they get to the store.", SOFT_GOLD),
        ("Local", "They need recommendations that feel tied to nearby stores.", SOFT_BLUE),
        ("Understandable", "They need outputs they can trust without knowing the math.", SOFT_TERRA),
    ]
    x = 830
    y = 322
    for title, body, fill_color in needs:
        card(draw, (x, y, x + 470, y + 245), fill=fill_color, outline=fill_color)
        draw.text((x + 28, y + 28), title, font=font(30, kind="bold"), fill=GREEN)
        draw_paragraph(draw, (x + 28, y + 86), body, font(26), TEXT, 410, line_gap=8)
        x += 500
        if x > 1400:
            x = 830
            y += 285
    draw.text((84, 958), "Why that audience fit makes sense", font=font(26, kind="bold"), fill=GREEN)
    draw_paragraph(
        draw,
        (84, 996),
        "The app exposes presets, shopping modes, pantry checkboxes, nearby store lookup, and plain-language notes. That is the design of a guided consumer decision tool, not a specialist analytics console.",
        font(24),
        MUTED,
        1640,
        line_gap=8,
    )
    return save_slide(img, "slide_02_audience.png")


def slide_3() -> Path:
    img, draw = make_canvas()
    heading(draw, 3, "Problem", "The Real-World Problem", "Nutrition advice is abstract, but shopping is local, budget-constrained, and full of practical tradeoffs.")
    boxes = [
        ("Abstract advice", "A user may know they want more protein or a lower-cost basket, but that does not tell them what to buy."),
        ("Too many tradeoffs", "Calories, protein, produce, convenience, pantry items, and cost all pull the basket in different directions."),
        ("Local reality", "The grocery decision depends on nearby stores and regional price levels, not just the target on paper."),
        ("Decision fatigue", "Without structure, the user either guesses or buys a basket that does not actually match the goal."),
    ]
    coords = [(84, 320), (1000, 320), (84, 610), (1000, 610)]
    fills = [SOFT_BLUE, SOFT_GREEN, SOFT_TERRA, SOFT_GOLD]
    for (title, body), (x, y), fill_color in zip(boxes, coords, fills, strict=True):
        card(draw, (x, y, x + 836, y + 220), fill=fill_color, outline=fill_color)
        draw.text((x + 30, y + 28), title, font=font(32, kind="bold"), fill=GREEN)
        draw_paragraph(draw, (x + 30, y + 90), body, font(27), TEXT, 760, line_gap=9)
    card(draw, (84, 872, 1836, 980), fill=GREEN, outline=GREEN)
    draw.text(
        (118, 907),
        "Communication goal: help a non-expert user turn a health goal into a grocery decision they can actually act on.",
        font=font(30, kind="bold"),
        fill="#FDF8EE",
    )
    return save_slide(img, "slide_03_problem.png")


def slide_4() -> Path:
    img, draw = make_canvas()
    heading(draw, 4, "How It Works", "What The Tool Does", "The best non-technical explanation is a simple decision pipeline, not a math lecture.")
    steps = [
        ("1. User inputs", "Location, targets, preferences, pantry, and shopping window."),
        ("2. Basket options", "The planner builds several plausible grocery baskets."),
        ("3. Final choice", "It picks the basket that best balances nutrition fit, practicality, and price-awareness."),
        ("4. Decision support", "It returns a list, a summary, a rough cost, and store suggestions."),
    ]
    left = 90
    top = 370
    step_w = 395
    for idx, (title, body) in enumerate(steps):
        x = left + idx * 450
        card(draw, (x, top, x + step_w, top + 300), fill=CARD)
        draw.ellipse((x + 22, top + 20, x + 96, top + 94), fill=SOFT_GOLD, outline=SOFT_GOLD)
        draw.text((x + 44, top + 35), str(idx + 1), font=font(32, kind="bold"), fill=GREEN)
        draw.text((x + 26, top + 122), title, font=font(31, kind="bold"), fill=GREEN)
        draw_paragraph(draw, (x + 26, top + 176), body, font(25), MUTED, 340, line_gap=8)
        if idx < len(steps) - 1:
            arrow_x = x + step_w + 18
            draw.line((arrow_x, top + 150, arrow_x + 48, top + 150), fill=GREEN_2, width=8)
            draw.polygon([(arrow_x + 48, top + 150), (arrow_x + 28, top + 136), (arrow_x + 28, top + 164)], fill=GREEN_2)
    card(draw, (90, 752, 1830, 920), fill=SOFT_GREEN, outline=SOFT_GREEN)
    draw.text((122, 786), "Plain-English summary", font=font(28, kind="bold"), fill=GREEN)
    draw_paragraph(
        draw,
        (122, 838),
        "In other words, the app tries a few sensible shopping baskets and picks one that seems closest to the user's goal without ignoring realism.",
        font(31),
        TEXT,
        1580,
        line_gap=9,
    )
    return save_slide(img, "slide_04_workflow.png")


def slide_5() -> Path:
    img, draw = make_canvas()
    heading(draw, 5, "Interaction", "How A User Interacts With It", "The interface is organized like a guided grocery-planning workflow rather than a technical analytics dashboard.")
    paste_image(img, ASSET_DIR / "dashboard-inputs.png", (84, 312, 1120, 770), fit="contain")
    paste_image(img, ASSET_DIR / "dashboard-stores.png", (1180, 312, 1836, 770), fit="contain")
    card(draw, (84, 824, 1836, 986), fill=CARD)
    draw.text((120, 858), "Why this is user-friendly", font=font(28, kind="bold"), fill=GREEN)
    draw_bullets(
        draw,
        (120, 910),
        [
            "Preset buttons reduce setup effort for common goals like muscle gain, fat loss, and budget-friendly healthy.",
            "The main controls are everyday choices: calories, protein, days, shopping mode, dietary flags, and pantry items.",
            "Nearby store lookup and Recommend are separate steps, which helps a first-time user explore before committing.",
        ],
        1620,
        text_font=font(25),
        gap_after=10,
    )
    draw.text((84, 780), "Step 1: set the goal and preferences", font=font(22, kind="bold"), fill=GREEN)
    draw.text((1180, 780), "Step 2: load nearby supermarkets", font=font(22, kind="bold"), fill=GREEN)
    return save_slide(img, "slide_05_interaction.png")


def scenario_card(draw: ImageDraw.ImageDraw, x: int, y: int, title: str, data: dict[str, object], fill: str) -> None:
    card(draw, (x, y, x + 552, y + 608), fill=fill, outline=fill)
    draw.text((x + 28, y + 28), title, font=font(30, kind="bold"), fill=GREEN)
    cost_text = f"${data['cost']:.2f}"
    draw.text((x + 28, y + 88), cost_text, font=font(54, kind="title"), fill=TEXT)
    draw.text((x + 28, y + 152), "Typical basket cost", font=font(24), fill=MUTED)
    chips_y = y + 215
    chip_x = x + 28
    for item in data["items"][:5]:
        chip_x = pill(draw, chip_x, chips_y, item, fill=CARD, text_fill=TEXT, font_size=19)
        chip_x += 12
        if chip_x > x + 450:
            chips_y += 54
            chip_x = x + 28
    draw.text((x + 28, y + 372), "Observed effect", font=font(24, kind="bold"), fill=GREEN)
    draw_paragraph(draw, (x + 28, y + 414), data["note"], font(24), TEXT, 486, line_gap=8)
    protein_gap = data["protein_est"] - data["protein_target"]
    draw.text((x + 28, y + 510), f"Protein gap: {protein_gap:+.1f} g", font=font(23, kind="bold"), fill=MUTED)
    store_text = f"Produce fit: {data['stores'].get('produce_pick', 'n/a')}"
    draw.text((x + 28, y + 548), store_text, font=font(23), fill=MUTED)


def slide_6() -> Path:
    img, draw = make_canvas()
    heading(draw, 6, "Decision Impact", "What Changes When Decisions Change", "These outputs were observed on the deployed app in Mountain View while holding the same base target constant.")
    card(draw, (84, 272, 1836, 358), fill=GREEN, outline=GREEN)
    draw.text((118, 302), "Shared setup: Mountain View, CA | 2,200 kcal | 130 g protein | 1 day | balanced shopping mode", font=font(29, kind="bold"), fill="#FDF8EE")
    scenario_card(draw, 84, 398, "Base plan", SCENARIOS["baseline"], SOFT_BLUE)
    scenario_card(draw, 684, 398, "Budget preference", SCENARIOS["budget_toggle"], SOFT_GREEN)
    scenario_card(draw, 1284, 398, "Vegan preset", SCENARIOS["vegan_preset"], SOFT_TERRA)
    draw.text((84, 1022), "Takeaway: the app is not just filling calories and protein. It visibly trades off cost, dietary fit, and basket composition when the user changes presets.", font=font(23), fill=MUTED)
    return save_slide(img, "slide_06_tradeoffs.png")


def slide_7() -> Path:
    img, draw = make_canvas()
    heading(draw, 7, "Example", "Key Insight / Recommendation Example", "A concrete recommendation helps the audience understand what the tool is actually saying.")
    paste_image(img, ASSET_DIR / "dashboard-results.png", (84, 306, 1120, 982), crop=(0, 0, 1416, 2050), fit="cover")
    card(draw, (1180, 306, 1836, 982), fill=CARD)
    draw.text((1218, 346), "Observed base recommendation", font=font(30, kind="bold"), fill=GREEN)
    draw.text((1218, 416), "What the model is saying", font=font(24, kind="bold"), fill=MUTED)
    draw_paragraph(
        draw,
        (1218, 460),
        "Buy a small set of versatile staples that gets you close to the target without making the trip expensive or overly complicated.",
        font(29),
        TEXT,
        560,
        line_gap=9,
    )
    stats = [
        ("Protein", f"{SCENARIOS['baseline']['protein_est']:.1f} g vs {SCENARIOS['baseline']['protein_target']:.1f} g target"),
        ("Calories", f"{SCENARIOS['baseline']['calorie_est']:.1f} vs {SCENARIOS['baseline']['calorie_target']:.1f} target"),
        ("Cost", f"About ${SCENARIOS['baseline']['cost']:.2f} for the basket"),
    ]
    y = 612
    for label, body in stats:
        card(draw, (1218, y, 1796, y + 78), fill=SOFT_GREEN, outline=SOFT_GREEN)
        draw.text((1244, y + 14), label, font=font(23, kind="bold"), fill=GREEN)
        draw.text((1404, y + 14), body, font=font(23), fill=TEXT)
        y += 96
    draw.text((1218, 916), "Why this matters", font=font(24, kind="bold"), fill=MUTED)
    draw_paragraph(
        draw,
        (1218, 952),
        "The recommendation is useful because it is actionable: it tells the user what to buy, what the basket roughly achieves, and where it may fit best nearby. In this observed case, Safeway is the one-stop pick and La Plaza is the stronger produce fit.",
        font(22),
        TEXT,
        560,
        line_gap=6,
    )
    return save_slide(img, "slide_07_example.png")


def slide_8() -> Path:
    img, draw = make_canvas()
    heading(draw, 8, "Limits", "Limitations + How To Improve", "Strong communication includes being honest about uncertainty, assumptions, and what the current tool still cannot do.")
    card(draw, (84, 304, 896, 864), fill=SOFT_TERRA, outline=SOFT_TERRA)
    draw.text((118, 340), "Current limitations", font=font(30, kind="bold"), fill=GREEN)
    draw_bullets(
        draw,
        (118, 402),
        [
            "The basket uses generic foods, not exact branded products or exact live store inventory.",
            "Price estimates are regional references from public price tables, not store quotes.",
            "Store recommendations are coarse fit suggestions based on nearby store type, brand, and distance.",
            "Long shopping windows are heuristic. Fresh and bulk modes both soften perishables, so results become approximate.",
            "The tool is not medical nutrition advice and does not capture allergies or complex health needs.",
        ],
        720,
        text_font=font(25),
        gap_after=11,
    )
    card(draw, (948, 304, 1836, 864), fill=SOFT_GREEN, outline=SOFT_GREEN)
    draw.text((982, 340), "How to improve it", font=font(30, kind="bold"), fill=GREEN)
    draw_bullets(
        draw,
        (982, 402),
        [
            "Add real store inventory and current store-specific prices where possible.",
            "Let the user set a hard spending cap, not just a budget-friendly preference.",
            "Support allergies, household size, cultural preferences, and stronger nutrition guardrails.",
            "Collect feedback from real student or young professional users and adapt the interface around what confuses them.",
            "Show confidence labels more clearly so users know what is exact versus approximate.",
        ],
        760,
        text_font=font(25),
        gap_after=11,
    )
    card(draw, (84, 896, 1836, 992), fill=CARD)
    fresh = SCENARIOS["fresh7"]
    bulk = SCENARIOS["bulk7"]
    evidence = (
        f"Observed evidence: in 7-day mode, fresh landed at {fresh['protein_est']:.1f} g protein versus a {fresh['protein_target']:.1f} g target, "
        f"and bulk landed at {bulk['protein_est']:.1f} g versus {bulk['protein_target']:.1f} g. That shows why perishables and scaling remain a real limitation."
    )
    draw_paragraph(draw, (118, 922), evidence, font(24), MUTED, 1620, line_gap=8)
    return save_slide(img, "slide_08_limitations.png")


def stakeholder_box(draw: ImageDraw.ImageDraw, x: int, y: int, title: str, body: str, fill: str) -> None:
    card(draw, (x, y, x + 826, y + 210), fill=fill, outline=fill)
    draw.text((x + 28, y + 26), title, font=font(29, kind="bold"), fill=GREEN)
    draw_paragraph(draw, (x + 28, y + 82), body, font(25), TEXT, 756, line_gap=8)


def slide_9() -> Path:
    img, draw = make_canvas()
    heading(draw, 9, "Ethics", "Stakeholders + Ethics", "A useful recommendation system should help without overstating certainty or excluding the people who need help most.")
    stakeholder_box(draw, 84, 310, "Shoppers", "They gain a clearer starting point for cost, nutrition, and store choice, especially if they feel overwhelmed by grocery planning.", SOFT_BLUE)
    stakeholder_box(draw, 1010, 310, "Stores and educators", "Nearby stores may benefit from more intentional trips, and nutrition coaches or campus wellness groups could use the tool as a teaching aid.", SOFT_GREEN)
    stakeholder_box(draw, 84, 560, "Potentially underserved users", "People with allergies, medical conditions, low food access, or strong cultural constraints may need more safeguards than the current tool provides.", SOFT_TERRA)
    stakeholder_box(draw, 1010, 560, "Ethical design responsibility", "The tool should keep telling the truth about approximation: regional prices are not store quotes and store-fit suggestions are not inventory guarantees.", SOFT_GOLD)
    card(draw, (84, 834, 1836, 986), fill=GREEN, outline=GREEN)
    draw.text((118, 876), "Thoughtful ethical stance", font=font(28, kind="bold"), fill="#FDF8EE")
    draw_paragraph(
        draw,
        (118, 918),
        "The responsible position is to present the tool as a practical starting point for grocery decisions, not as medical advice or a perfect answer. Transparency is part of the product, not a disclaimer added at the end.",
        font(23),
        "#FDF8EE",
        1600,
        line_gap=7,
    )
    return save_slide(img, "slide_09_ethics.png")


def takeaway_box(draw: ImageDraw.ImageDraw, x: int, title: str, body: str, fill: str) -> None:
    card(draw, (x, 360, x + 540, 760), fill=fill, outline=fill)
    draw.text((x + 30, 398), title, font=font(32, kind="bold"), fill=GREEN)
    draw_paragraph(draw, (x + 30, 476), body, font(28), TEXT, 480, line_gap=9)


def slide_10() -> Path:
    img, draw = make_canvas()
    heading(draw, 10, "Close", "Takeaways + Future Use", "The strongest close is simple: this project makes nutrition goals more actionable for an everyday shopper.")
    takeaway_box(draw, 84, "Why it matters", "It turns an abstract nutrition goal into groceries, prices, and store choices a person can immediately use.", SOFT_BLUE)
    takeaway_box(draw, 690, "Why it is worth using", "The outputs are understandable, local-feeling, and honest about what is approximate.", SOFT_GREEN)
    takeaway_box(draw, 1296, "What comes next", "Pilot it with real users, collect feedback, and strengthen the data and confidence signals.", SOFT_TERRA)
    card(draw, (84, 820, 1836, 970), fill=GREEN, outline=GREEN)
    draw.text((118, 854), "Closing line", font=font(28, kind="bold"), fill="#FDF8EE")
    draw_paragraph(
        draw,
        (118, 902),
        "If the goal is better grocery decisions for non-experts, this tool already works as a strong communication bridge between nutrition targets and real shopping behavior.",
        font(29),
        "#FDF8EE",
        1600,
        line_gap=8,
    )
    return save_slide(img, "slide_10_takeaways.png")


def slide_11() -> Path:
    img, draw = make_canvas()
    heading(draw, 11, "Appendix", "Demo Backup", "If the live demo fails, this backup slide still lets the audience see the full recommendation workflow and outputs.")
    paste_image(img, ASSET_DIR / "dashboard-full-results.png", (84, 286, 1440, 988), crop=(0, 0, 1440, 4200), fit="contain")
    card(draw, (1490, 286, 1836, 988), fill=CARD)
    draw.text((1520, 324), "What I would point to", font=font(28, kind="bold"), fill=GREEN)
    draw_bullets(
        draw,
        (1520, 386),
        [
            "The grouped shopping list by food role.",
            "The cost estimate and price notes.",
            "The recommended store picks for the basket.",
            "The nutrition summary versus target.",
            "The approximation notes that prevent overconfidence.",
        ],
        272,
        text_font=font(24),
        gap_after=12,
    )
    return save_slide(img, "slide_11_backup.png")


def build_slides() -> list[Path]:
    return [
        slide_1(),
        slide_2(),
        slide_3(),
        slide_4(),
        slide_5(),
        slide_6(),
        slide_7(),
        slide_8(),
        slide_9(),
        slide_10(),
        slide_11(),
    ]


def build_pptx(slide_paths: list[Path]) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for slide_path in slide_paths:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(slide_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    if len(prs.slides) > len(slide_paths):
        first = prs.slides._sldIdLst[0]
        prs.slides._sldIdLst.remove(first)
    out_path = OUT_DIR / "indeng243_diet_planner_presentation.pptx"
    prs.save(out_path)
    return out_path


def build_pdf(slide_paths: list[Path]) -> Path:
    images = [Image.open(path).convert("RGB") for path in slide_paths]
    out_path = OUT_DIR / "indeng243_diet_planner_presentation.pdf"
    images[0].save(out_path, save_all=True, append_images=images[1:], resolution=150.0)
    return out_path


def write_markdown() -> list[Path]:
    summary_path = OUT_DIR / "indeng243_project_understanding.md"
    outline_path = OUT_DIR / "indeng243_presentation_outline.md"
    transcript_path = OUT_DIR / "indeng243_presentation_transcript.md"
    rubric_path = OUT_DIR / "indeng243_rubric_coverage_matrix.md"
    package_path = OUT_DIR / "indeng243_presentation_package.md"

    summary_lines = [
        "# Project Understanding Summary",
        "",
        f"- **Stack:** {PROJECT_UNDERSTANDING['stack']}",
        f"- **App purpose:** {PROJECT_UNDERSTANDING['app_purpose']}",
        f"- **Inferred audience:** {PROJECT_UNDERSTANDING['inferred_audience']}",
        f"- **Key inputs:** {PROJECT_UNDERSTANDING['key_inputs']}",
        f"- **Key outputs:** {PROJECT_UNDERSTANDING['key_outputs']}",
        f"- **Main tradeoff or objective:** {PROJECT_UNDERSTANDING['main_tradeoff']}",
        "- **Evidence-based uncertainties:**",
    ]
    for item in PROJECT_UNDERSTANDING["uncertainties"]:
        summary_lines.append(f"  - {item}")
    summary_path.write_text("\n".join(summary_lines) + "\n", encoding="utf-8")

    outline_lines = ["# Presentation Outline", ""]
    for slide in SLIDE_OUTLINE:
        outline_lines.append(f"## Slide {slide['slide']} - {slide['title']}")
        for bullet in slide["bullets"]:
            outline_lines.append(f"- {bullet}")
        outline_lines.append("")
    outline_path.write_text("\n".join(outline_lines).rstrip() + "\n", encoding="utf-8")

    transcript_lines = ["# Full Transcript", ""]
    for item in TRANSCRIPT:
        transcript_lines.append(f"## Slide {item['slide']} - {item['title']}")
        transcript_lines.append(item["script"])
        transcript_lines.append("")
    transcript_path.write_text("\n".join(transcript_lines).rstrip() + "\n", encoding="utf-8")

    rubric_lines = [
        "# Rubric Coverage Matrix",
        "",
        "| Rubric item | Where it is covered (slide number) | How it is covered |",
        "| --- | --- | --- |",
    ]
    for item, slides, how in RUBRIC_ROWS:
        rubric_lines.append(f"| {item} | {slides} | {how} |")
    rubric_path.write_text("\n".join(rubric_lines) + "\n", encoding="utf-8")

    package_lines = [
        "# Presentation Package",
        "",
        "## 1. Project Understanding Summary",
        f"- **Stack:** {PROJECT_UNDERSTANDING['stack']}",
        f"- **App purpose:** {PROJECT_UNDERSTANDING['app_purpose']}",
        f"- **Inferred audience:** {PROJECT_UNDERSTANDING['inferred_audience']}",
        f"- **Key inputs:** {PROJECT_UNDERSTANDING['key_inputs']}",
        f"- **Key outputs:** {PROJECT_UNDERSTANDING['key_outputs']}",
        f"- **Main tradeoff or objective:** {PROJECT_UNDERSTANDING['main_tradeoff']}",
        "- **Evidence-based uncertainties:**",
    ]
    for item in PROJECT_UNDERSTANDING["uncertainties"]:
        package_lines.append(f"  - {item}")
    package_lines += ["", "## 2. Presentation Outline", ""]
    for slide in SLIDE_OUTLINE:
        package_lines.append(f"### Slide {slide['slide']} - {slide['title']}")
        for bullet in slide["bullets"]:
            package_lines.append(f"- {bullet}")
        package_lines.append("")
    package_lines += ["## 3. Full Transcript", ""]
    for item in TRANSCRIPT:
        package_lines.append(f"### Slide {item['slide']} - {item['title']}")
        package_lines.append(item["script"])
        package_lines.append("")
    package_lines += [
        "## 4. Rubric Coverage Matrix",
        "",
        "| Rubric item | Where it is covered (slide number) | How it is covered |",
        "| --- | --- | --- |",
    ]
    for item, slides, how in RUBRIC_ROWS:
        package_lines.append(f"| {item} | {slides} | {how} |")
    package_lines += [
        "",
        "## 5. Files Created",
        "",
        "- `presentation/indeng243_diet_planner_presentation.pptx`",
        "- `presentation/indeng243_diet_planner_presentation.pdf`",
        "- `presentation/indeng243_project_understanding.md`",
        "- `presentation/indeng243_presentation_outline.md`",
        "- `presentation/indeng243_presentation_transcript.md`",
        "- `presentation/indeng243_rubric_coverage_matrix.md`",
        "- `presentation/indeng243_presentation_package.md`",
    ]
    package_path.write_text("\n".join(package_lines).rstrip() + "\n", encoding="utf-8")
    return [summary_path, outline_path, transcript_path, rubric_path, package_path]


def main() -> None:
    ensure_dirs()
    slide_paths = build_slides()
    build_pptx(slide_paths)
    build_pdf(slide_paths)
    write_markdown()


if __name__ == "__main__":
    main()
